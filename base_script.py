"""
FreightAI Architecture Presentation Generator
Creates a PowerPoint presentation with architecture diagrams and documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def create_section_slide(prs, title):
    """Create a section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section Header layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style the section title
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    
    return slide

def create_content_slide(prs, title):
    """Create a content slide with title"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """Add a text box to a slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    
    return textbox

def add_shape_box(slide, left, top, width, height, text, fill_color=None, border_color=None):
    """Add a rectangular shape with text"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    # Set fill color
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 240, 250)
    
    # Set border
    shape.line.color.rgb = border_color if border_color else RGBColor(0, 102, 204)
    shape.line.width = Pt(2)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    
    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """Add an arrow connector between two points"""
    connector = slide.shapes.add_connector(
        2,  # Straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(100, 100, 100)
    connector.line.width = Pt(2)
    return connector

def create_overview_slide(prs):
    """Create project overview slide"""
    slide = create_content_slide(prs, "FreightAI - Project Overview")
    
    # Add overview text
    content = """The FreightAI application is a comprehensive logistics management system that handles:

• Vehicle Inspections - Automated inspection processing with AI-powered analysis
• POD (Proof of Delivery) Management - Document extraction and storage
• Image Management - Secure image upload and retrieval with presigned URLs
• Real-time Data Access - Fast, scalable data retrieval through DynamoDB

Built on AWS serverless architecture for high availability and cost efficiency."""
    
    add_textbox(slide, 0.5, 1.5, 9, 3, content, font_size=18)
    
    # Add key features box
    add_shape_box(slide, 0.5, 4.8, 4, 1.5, 
                  "Key Technologies\n\nReact • TypeScript\nAWS Lambda • API Gateway\nDynamoDB • S3",
                  fill_color=RGBColor(220, 235, 255))
    
    add_shape_box(slide, 5.5, 4.8, 4, 1.5,
                  "AWS Account\n\nAccount: 813281204422\nRegion: us-east-1",
                  fill_color=RGBColor(255, 240, 220))
    
    return slide

def create_architecture_diagram(prs):
    """Create high-level architecture diagram"""
    slide = create_content_slide(prs, "High-Level Architecture")
    
    # Frontend
    add_shape_box(slide, 3.5, 1.3, 3, 0.8, "React Frontend\nTypeScript",
                  fill_color=RGBColor(97, 218, 251))
    
    # API Gateway
    add_shape_box(slide, 3.5, 2.5, 3, 0.8, "API Gateway\nmc7upldn34",
                  fill_color=RGBColor(255, 153, 51))
    
    # Lambda react-ui
    add_shape_box(slide, 3.5, 3.7, 3, 0.8, "Lambda: react-ui\nMain API Handler",
                  fill_color=RGBColor(255, 153, 51))
    
    # DynamoDB left
    add_shape_box(slide, 0.5, 5.2, 2, 0.8, "DynamoDB\nPOD Table",
                  fill_color=RGBColor(51, 102, 255))
    
    # DynamoDB right
    add_shape_box(slide, 2.8, 5.2, 2.3, 0.8, "DynamoDB\ninspection_results",
                  fill_color=RGBColor(51, 102, 255))
    
    # Lambda functions
    add_shape_box(slide, 5.5, 5.2, 1.8, 0.8, "Lambda\npod-extract2",
                  fill_color=RGBColor(255, 153, 51))
    
    add_shape_box(slide, 7.5, 5.2, 1.8, 0.8, "Lambda\nimage-presign",
                  fill_color=RGBColor(255, 153, 51))
    
    # S3 bucket
    add_shape_box(slide, 6.5, 6.3, 2.8, 0.8, "S3 Buckets\nImages & POD Storage",
                  fill_color=RGBColor(76, 175, 80))
    
    # Add arrows
    add_arrow(slide, 5, 2.1, 5, 2.5)  # Frontend to API Gateway
    add_arrow(slide, 5, 3.3, 5, 3.7)  # API Gateway to Lambda
    add_arrow(slide, 3.8, 4.5, 1.5, 5.2)  # Lambda to POD table
    add_arrow(slide, 4.5, 4.5, 3.9, 5.2)  # Lambda to inspection_results
    add_arrow(slide, 5.5, 4.5, 6.4, 5.2)  # Lambda to pod-extract2
    add_arrow(slide, 6, 4.5, 8.4, 5.2)  # Lambda to image-presign
    add_arrow(slide, 7.9, 6.0, 7.9, 6.3)  # Lambda to S3
    
    return slide

def create_api_gateway_slide(prs):
    """Create API Gateway details slide"""
    slide = create_content_slide(prs, "API Gateway Configuration")
    
    # API details box
    add_shape_box(slide, 0.5, 1.3, 4.5, 1.2,
                  "HTTP API: react-api\n" +
                  "ID: mc7upldn34\n" +
                  "Region: us-east-1",
                  fill_color=RGBColor(255, 153, 51))
    
    # Endpoint box
    add_textbox(slide, 0.5, 2.7, 9, 0.4, 
                "Endpoint: https://mc7upldn34.execute-api.us-east-1.amazonaws.com",
                font_size=12, bold=True)
    
    # Routes
    routes_text = """API Routes:

• GET /ft1/inspections - Retrieve vehicle inspection records
• GET /pod - Retrieve POD documents
• ANY /presign - Generate presigned URLs for S3 access
• GET /s3-list - List S3 bucket contents
• ANY /ft1/presign - Generate presigned URLs (alternate endpoint)

Integration Type: AWS_PROXY
Target: Lambda function 'react-ui'
Timeout: 30 seconds"""
    
    add_textbox(slide, 0.5, 3.3, 9, 3, routes_text, font_size=14)
    
    return slide

def create_lambda_functions_slide(prs):
    """Create Lambda Functions overview slide"""
    slide = create_content_slide(prs, "Lambda Functions")
    
    # react-ui
    add_shape_box(slide, 0.5, 1.5, 4.3, 1.3,
                  "react-ui\n\n" +
                  "Main API handler for all routes\n" +
                  "Runtime: Python 3.11\n" +
                  "Integration: API Gateway",
                  fill_color=RGBColor(255, 200, 100))
    
    # pod-extract2
    add_shape_box(slide, 5.2, 1.5, 4.3, 1.3,
                  "pod-extract2\n\n" +
                  "POD document extraction\n" +
                  "Runtime: Python 3.11\n" +
                  "Process: Extract & store data",
                  fill_color=RGBColor(255, 200, 100))
    
    # image-presign
    add_shape_box(slide, 0.5, 3.2, 4.3, 1.3,
                  "image-presign\n\n" +
                  "Generate S3 presigned URLs\n" +
                  "Runtime: Python 3.11\n" +
                  "Purpose: Secure image access",
                  fill_color=RGBColor(255, 200, 100))
    
    # image-load
    add_shape_box(slide, 5.2, 3.2, 4.3, 1.3,
                  "image-load\n\n" +
                  "Load and process S3 images\n" +
                  "Runtime: Python 3.11\n" +
                  "Purpose: Image management",
                  fill_color=RGBColor(255, 200, 100))
    
    # Common config
    config_text = "Common Configuration:\n• Handler: lambda_function.lambda_handler\n• Region: us-east-1\n• Account: 813281204422"
    add_textbox(slide, 0.5, 4.8, 9, 1, config_text, font_size=14, bold=True)
    
    return slide

def create_dynamodb_slide(prs):
    """Create DynamoDB tables slide"""
    slide = create_content_slide(prs, "DynamoDB Tables")
    
    # POD Table
    add_shape_box(slide, 0.5, 1.5, 4.3, 2.2,
                  "POD Table\n\n" +
                  "Partition Key: PODID (String)\n" +
                  "Sort Key: Delivery_date (String)\n\n" +
                  "Billing: PAY_PER_REQUEST\n" +
                  "Items: ~10\n" +
                  "Read: 12,000/sec\n" +
                  "Write: 4,000/sec",
                  fill_color=RGBColor(100, 150, 255))
    
    # inspection_results Table
    add_shape_box(slide, 5.2, 1.5, 4.3, 2.2,
                  "inspection_results\n\n" +
                  "Partition Key: inspectionId\n" +
                  "Sort Key: license_plate\n\n" +
                  "Billing: PAY_PER_REQUEST\n" +
                  "Items: ~2\n" +
                  "Read: 12,000/sec\n" +
                  "Write: 4,000/sec",
                  fill_color=RGBColor(100, 150, 255))
    
    # Common features
    features_text = """Key Features:

• On-Demand billing for cost optimization
• STANDARD table class for consistent performance
• Warm throughput enabled for low-latency access
• Region: us-east-1 (US East - N. Virginia)
• Deletion protection: Disabled (development environment)"""
    
    add_textbox(slide, 0.5, 4, 9, 2, features_text, font_size=14)
    
    return slide

def create_s3_storage_slide(prs):
    """Create S3 storage slide"""
    slide = create_content_slide(prs, "S3 Storage Architecture")
    
    # Inspection images bucket
    add_shape_box(slide, 0.5, 1.5, 4.3, 1.8,
                  "aws-s3-futuregen-images\n\n" +
                  "Vehicle Inspection Images\n\n" +
                  "Folders:\n" +
                  "• images-pass/ (successful)\n" +
                  "• images-fail/ (failed)",
                  fill_color=RGBColor(100, 200, 100))
    
    # POD documents bucket
    add_shape_box(slide, 5.2, 1.5, 4.3, 1.8,
                  "aws-s3-futuregen-pod1\n\n" +
                  "POD Documents\n\n" +
                  "Storage:\n" +
                  "• Delivery receipts\n" +
                  "• Signed documents",
                  fill_color=RGBColor(100, 200, 100))
    
    # CORS configuration
    cors_text = """CORS Configuration:

• Allow Origins: * (all origins)
• Allow Methods: GET, POST, PUT, DELETE
• Allow Headers: content-type, authorization, x-api-key, x-amz-date, x-amz-security-token
• Max Age: 300 seconds

Security: Presigned URLs with 6-hour expiration for secure access"""
    
    add_textbox(slide, 0.5, 3.7, 9, 2.5, cors_text, font_size=14)
    
    return slide

def create_data_flow_diagram(prs):
    """Create data flow diagram"""
    slide = create_content_slide(prs, "Data Flow - Vehicle Inspections")
    
    # Step 1
    add_shape_box(slide, 0.3, 1.5, 1.8, 0.8, "1. User Request\nGET /inspections",
                  fill_color=RGBColor(200, 230, 255))
    
    # Step 2
    add_shape_box(slide, 2.4, 1.5, 1.8, 0.8, "2. API Gateway\nRoute Request",
                  fill_color=RGBColor(255, 200, 150))
    
    # Step 3
    add_shape_box(slide, 4.5, 1.5, 1.8, 0.8, "3. Lambda\nreact-ui",
                  fill_color=RGBColor(255, 200, 150))
    
    # Step 4
    add_shape_box(slide, 6.6, 1.5, 1.8, 0.8, "4. Query\nDynamoDB",
                  fill_color=RGBColor(150, 180, 255))
    
    # Step 5
    add_shape_box(slide, 1.5, 2.8, 2.2, 0.8, "5. Generate\nPresigned URLs",
                  fill_color=RGBColor(255, 200, 150))
    
    # Step 6
    add_shape_box(slide, 4.2, 2.8, 2.2, 0.8, "6. Return JSON\nwith ImageURLs",
                  fill_color=RGBColor(200, 255, 200))
    
    # Step 7
    add_shape_box(slide, 6.9, 2.8, 2.2, 0.8, "7. Frontend\nDisplay Data",
                  fill_color=RGBColor(200, 230, 255))
    
    # Add arrows
    add_arrow(slide, 2.1, 1.9, 2.4, 1.9)
    add_arrow(slide, 4.2, 1.9, 4.5, 1.9)
    add_arrow(slide, 6.3, 1.9, 6.6, 1.9)
    add_arrow(slide, 7.5, 2.3, 2.6, 2.8)
    add_arrow(slide, 3.7, 3.2, 4.2, 3.2)
    add_arrow(slide, 6.4, 3.2, 6.9, 3.2)
    
    # Add details
    details = """Process Details:

• Inspections are stored in DynamoDB with S3 image references
• Lambda generates presigned URLs (6-hour expiry) for secure image access
• Response includes inspection metadata, scores, and direct image URLs
• Frontend displays inspection results with images"""
    
    add_textbox(slide, 0.5, 4.2, 9, 2, details, font_size=14)
    
    return slide

def create_deployment_slide(prs):
    """Create deployment architecture slide"""
    slide = create_content_slide(prs, "Deployment Architecture")
    
    content = """Deployment Strategy:

Frontend Deployment:
• React application built with TypeScript
• Build command: npm run build
• Output: Optimized static files in /build directory
• Deployment: Can be hosted on S3 + CloudFront, Amplify, or any static hosting

Backend Deployment:
• Lambda functions deployed via AWS CLI or console
• Functions packaged as .zip files with dependencies
• API Gateway routes traffic to Lambda functions
• DynamoDB tables pre-created with on-demand billing

Environment Configuration:
• API endpoint configured in React environment
• Lambda environment variables set per function
• CORS enabled for cross-origin requests
• S3 bucket policies configured for Lambda access"""
    
    add_textbox(slide, 0.5, 1.5, 9, 4.5, content, font_size=14)
    
    return slide

def create_security_slide(prs):
    """Create security considerations slide"""
    slide = create_content_slide(prs, "Security & Best Practices")
    
    # Current security
    add_shape_box(slide, 0.5, 1.5, 4.3, 2.3,
                  "Current Security\n\n" +
                  "✓ Presigned URLs (6h expiry)\n" +
                  "✓ IAM roles for Lambda\n" +
                  "✓ HTTPS endpoints\n" +
                  "✓ CORS configuration\n" +
                  "✓ On-demand scaling",
                  fill_color=RGBColor(200, 255, 200))
    
    # Recommended enhancements
    add_shape_box(slide, 5.2, 1.5, 4.3, 2.3,
                  "Recommended for Production\n\n" +
                  "• API key authentication\n" +
                  "• WAF (Web Application Firewall)\n" +
                  "• CloudWatch monitoring\n" +
                  "• DynamoDB backups\n" +
                  "• Restrict CORS origins",
                  fill_color=RGBColor(255, 240, 200))
    
    # Cost optimization
    cost_text = """Cost Optimization:

• DynamoDB on-demand pricing - pay per request
• Lambda charged per invocation and execution time
• S3 storage costs based on data volume
• API Gateway costs per API call
• Estimated monthly cost: $5-50 depending on usage"""
    
    add_textbox(slide, 0.5, 4.2, 9, 2, cost_text, font_size=14)
    
    return slide

def create_api_endpoints_slide(prs):
    """Create API endpoints detail slide"""
    slide = create_content_slide(prs, "API Endpoints Detail")
    
    endpoints = """Main API Endpoints:

1. GET /ft1/inspections
   • Returns all vehicle inspection records
   • Includes presigned image URLs (6-hour expiry)
   • Response includes inspection scores and status

2. GET /pod
   • Retrieves POD (Proof of Delivery) documents
   • Returns document metadata and S3 locations

3. ANY /presign or /ft1/presign
   • Generates presigned URLs for S3 objects
   • Allows secure, temporary access to private files

4. GET /s3-list
   • Lists contents of S3 buckets
   • Returns available images and documents"""
    
    add_textbox(slide, 0.5, 1.5, 9, 4.5, endpoints, font_size=14)
    
    return slide

def create_future_enhancements_slide(prs):
    """Create future enhancements slide"""
    slide = create_content_slide(prs, "Future Enhancements")
    
    # Phase 1
    add_shape_box(slide, 0.5, 1.5, 4.3, 1.5,
                  "Phase 1: Security\n\n" +
                  "• API authentication\n" +
                  "• User management\n" +
                  "• Audit logging",
                  fill_color=RGBColor(255, 220, 220))
    
    # Phase 2
    add_shape_box(slide, 5.2, 1.5, 4.3, 1.5,
                  "Phase 2: Features\n\n" +
                  "• Real-time notifications\n" +
                  "• Advanced AI analysis\n" +
                  "• Mobile app integration",
                  fill_color=RGBColor(220, 240, 255))
    
    # Phase 3
    add_shape_box(slide, 0.5, 3.3, 4.3, 1.5,
                  "Phase 3: Scale\n\n" +
                  "• Multi-region deployment\n" +
                  "• CDN integration\n" +
                  "• Performance optimization",
                  fill_color=RGBColor(220, 255, 220))
    
    # Phase 4
    add_shape_box(slide, 5.2, 3.3, 4.3, 1.5,
                  "Phase 4: Analytics\n\n" +
                  "• Business intelligence\n" +
                  "• Trend analysis\n" +
                  "• Reporting dashboards",
                  fill_color=RGBColor(255, 240, 200))
    
    return slide

def main():
    """Main function to create the presentation"""
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    create_title_slide(prs, "FreightAI", "AWS Architecture & System Design\nDecember 2025")
    create_overview_slide(prs)
    create_section_slide(prs, "System Architecture")
    create_architecture_diagram(prs)
    create_data_flow_diagram(prs)
    create_section_slide(prs, "AWS Components")
    create_api_gateway_slide(prs)
    create_lambda_functions_slide(prs)
    create_dynamodb_slide(prs)
    create_s3_storage_slide(prs)
    create_section_slide(prs, "API & Integration")
    create_api_endpoints_slide(prs)
    create_section_slide(prs, "Deployment & Operations")
    create_deployment_slide(prs)
    create_security_slide(prs)
    create_section_slide(prs, "Future Roadmap")
    create_future_enhancements_slide(prs)
    
    # Save presentation
    filename = "FreightAI_Architecture_Presentation.pptx"
    prs.save(filename)
    print(f"✓ Presentation created successfully: {filename}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    main()
